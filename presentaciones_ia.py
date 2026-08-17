"""
presentaciones_ia.py — Fábrica Visual IA Nivel Dios (v2)
═══════════════════════════════════════════════════════════════════════════
🖼️ IMÁGENES IA AUTOMÁTICAS (Pollinations.ai, sin API key):
   • Una imagen única por diapositiva en PowerPoint
   • Hero visual + imágenes por sección en infografías
   • Galería regenerable (cambia cualquier imagen con un clic)
🔍 AUTO-INSPECCIÓN INTELIGENTE:
   • La IA evalúa su propio trabajo (completitud, coherencia, calidad)
   • Semáforo visual: Excelente / Aceptable / Requiere mejoras
   • Métricas detalladas por dimensión
🎨 5 temas de diseño · 4 tonos · anclaje curricular (PDF/texto)
"""
import html as _html
import re
import urllib.parse
import random
import hashlib
from io import BytesIO
import streamlit as st
import streamlit.components.v1 as components
from core import ia

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor as PColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    PPTX_OK = True
except Exception:
    PPTX_OK = False

try:
    from pypdf import PdfReader
except Exception:
    try:
        from PyPDF2 import PdfReader
    except Exception:
        PdfReader = None

# ═══════════════════════════════════════════════════════════════════════════
# TEMAS DE DISEÑO
# ═══════════════════════════════════════════════════════════════════════════
TEMAS_PPT = {
    "🔵 Institucional Azul":  {"bg": "0F172A", "accent": "3B82F6", "text": "F8FAFC", "sub": "94A3B8"},
    "🌊 Moderno Turquesa":    {"bg": "F0FDFA", "accent": "0D9488", "text": "0F172A", "sub": "475569"},
    "🌙 Ejecutivo Oscuro":    {"bg": "111827", "accent": "F59E0B", "text": "F9FAFB", "sub": "9CA3AF"},
    "🔥 Vibrante Naranja":    {"bg": "FFF7ED", "accent": "EA580C", "text": "1C1917", "sub": "57534E"},
    "💜 Elegante Púrpura":    {"bg": "FAF5FF", "accent": "7C3AED", "text": "1E1B4B", "sub": "6B7280"},
}
TEMAS_INFO = {
    "🔵 Azul Institucional": {"bg": "#F0F4F8", "accent": "#2563EB", "accent2": "#3B82F6",
                               "accentsoft": "#DBEAFE", "card": "#FFFFFF", "text": "#0F172A", "muted": "#64748B"},
    "🌊 Turquesa Moderno":   {"bg": "#F0FDFA", "accent": "#0D9488", "accent2": "#14B8A6",
                               "accentsoft": "#CCFBF1", "card": "#FFFFFF", "text": "#134E4A", "muted": "#5B7C7A"},
    "🔥 Naranja Vibrante":   {"bg": "#FFF7ED", "accent": "#EA580C", "accent2": "#F97316",
                               "accentsoft": "#FFEDD5", "card": "#FFFFFF", "text": "#7C2D12", "muted": "#9A6A4F"},
    "💜 Púrpura Elegante":   {"bg": "#FAF5FF", "accent": "#7C3AED", "accent2": "#A78BFA",
                               "accentsoft": "#EDE9FE", "card": "#FFFFFF", "text": "#4C1D95", "muted": "#7B6D8F"},
    "🌙 Oscuro Ejecutivo":   {"bg": "#0F172A", "accent": "#F59E0B", "accent2": "#FBBF24",
                               "accentsoft": "#3B3424", "card": "#1E293B", "text": "#F8FAFC", "muted": "#94A3B8"},
}
TONOS = ["Formal académico", "Didáctico y claro", "Creativo e inspirador", "Técnico profesional"]

# ═══════════════════════════════════════════════════════════════════════════
# GENERADOR DE IMÁGENES IA (Pollinations.ai · sin API key)
# ═══════════════════════════════════════════════════════════════════════════
def generar_url_imagen(prompt: str, width: int = 1024, height: int = 576,
                       seed: int = None, estilo: str = "photorealistic") -> str:
    """Genera URL de imagen IA vía Pollinations.ai (gratis, sin API key)."""
    if seed is None:
        seed = int(hashlib.md5(prompt.encode()).hexdigest()[:8], 16) % 100000
    estilos = {
        "photorealistic": ", photorealistic, 4k, high quality, professional lighting",
        "illustration":   ", flat illustration, modern vector art, vibrant colors",
        "3d":             ", 3d render, octane render, ultra detailed, cinematic",
        "minimalist":     ", minimalist, clean design, simple composition",
        "academic":       ", educational illustration, clean, professional, diagram style",
    }
    suf = estilos.get(estilo, estilos["photorealistic"])
    full_prompt = f"{prompt}{suf}"
    encoded = urllib.parse.quote(full_prompt)
    return (f"https://image.pollinations.ai/prompt/{encoded}"
            f"?width={width}&height={height}&nologo=true&seed={seed}&model=flux")


def generar_imagenes_para_presentacion(datos: dict, tema_key: str) -> dict:
    """Genera una imagen por diapositiva + portada + conclusión."""
    tema = TEMAS_PPT.get(tema_key, TEMAS_PPT["🔵 Institucional Azul"])
    imagenes = {}
    titulo = datos.get("titulo", "Presentación educativa")

    # Portada
    imagenes["portada"] = generar_url_imagen(
        f"Conceptual abstract background for {titulo}, theme color #{tema['accent']}, "
        f"professional presentation cover, elegant composition",
        width=1280, height=720, estilo="minimalist"
    )

    # Una imagen por diapositiva de contenido
    for i, diap in enumerate(datos.get("diapositivas", []), 1):
        titulo_diap = diap.get("titulo", "")
        puntos = ", ".join(diap.get("puntos", [])[:2])
        prompt_img = (f"Educational visual for '{titulo_diap}', key concepts: {puntos}, "
                      f"professional teaching slide illustration, no text, clean composition")
        imagenes[f"diap_{i}"] = generar_url_imagen(
            prompt_img, width=800, height=500, seed=i * 17, estilo="academic"
        )

    # Conclusión
    imagenes["conclusion"] = generar_url_imagen(
        f"Closing visual for educational presentation about {titulo}, "
        f"success, achievement, inspirational concept, no text",
        width=1024, height=576, estilo="3d"
    )
    return imagenes


def generar_imagenes_para_infografia(datos: dict, tema_key: str) -> dict:
    """Genera hero + imágenes por sección para infografía."""
    imagenes = {}
    titulo = datos.get("titulo", "Infografía educativa")

    # Hero
    imagenes["hero"] = generar_url_imagen(
        f"Hero banner illustration for infographic about {titulo}, "
        f"vibrant, engaging, educational theme, no text, wide composition",
        width=1200, height=400, estilo="illustration"
    )

    # Una imagen por sección
    for i, sec in enumerate(datos.get("secciones", []), 1):
        titulo_sec = sec.get("titulo", "")
        desc_sec = sec.get("descripcion", "")[:80]
        imagenes[f"sec_{i}"] = generar_url_imagen(
            f"Illustration for '{titulo_sec}': {desc_sec}, educational infographic icon style, "
            f"clean, no text, square composition",
            width=400, height=400, seed=i * 31 + 7, estilo="illustration"
        )
    return imagenes

# ═══════════════════════════════════════════════════════════════════════════
# AUTO-INSPECCIÓN INTELIGENTE
# ═══════════════════════════════════════════════════════════════════════════
def inspeccionar_presentacion(datos: dict, num_esperado: int) -> dict:
    """Evalúa la calidad de la presentación generada. Devuelve métricas."""
    metricas = []
    total_puntos = 0
    max_puntos = 0

    # 1. Título (10 pts)
    titulo = datos.get("titulo", "")
    max_puntos += 10
    if titulo and len(titulo) >= 10 and len(titulo) <= 80:
        metricas.append(("📌 Título", 10, "Excelente", "Longitud ideal y claro"))
        total_puntos += 10
    elif titulo:
        metricas.append(("📌 Título", 6, "Aceptable", "Presente pero mejorable"))
        total_puntos += 6
    else:
        metricas.append(("📌 Título", 0, "Ausente", "Falta el título"))

    # 2. Subtítulo (5 pts)
    sub = datos.get("subtitulo", "")
    max_puntos += 5
    if sub and len(sub) >= 15:
        metricas.append(("🏷️ Subtítulo", 5, "Excelente", "Atractivo y descriptivo"))
        total_puntos += 5
    elif sub:
        metricas.append(("🏷️ Subtítulo", 3, "Aceptable", "Muy breve"))
        total_puntos += 3
    else:
        metricas.append(("🏷️ Subtítulo", 0, "Ausente", "Sin subtítulo"))

    # 3. Objetivo (15 pts)
    obj = datos.get("objetivo", "")
    max_puntos += 15
    if obj and len(obj) >= 40:
        tiene_verbo = any(v in obj.lower() for v in
            ["identificar", "analizar", "comprender", "aplicar", "evaluar",
             "crear", "desarrollar", "demostrar", "explicar", "distinguir"])
        if tiene_verbo:
            metricas.append(("🎯 Objetivo", 15, "Excelente", "Medible con verbo taxonómico"))
            total_puntos += 15
        else:
            metricas.append(("🎯 Objetivo", 10, "Bueno", "Claro pero sin verbo observable"))
            total_puntos += 10
    elif obj:
        metricas.append(("🎯 Objetivo", 5, "Aceptable", "Demasiado breve"))
        total_puntos += 5
    else:
        metricas.append(("🎯 Objetivo", 0, "Ausente", "Sin objetivo de aprendizaje"))

    # 4. Diapositivas (40 pts)
    diaps = datos.get("diapositivas", [])
    max_puntos += 40
    if len(diaps) == num_esperado:
        metricas.append(("📊 Cantidad de diapositivas", 10, "Excelente",
                         f"Exactamente {num_esperado} como solicitaste"))
        total_puntos += 10
    elif len(diaps) > 0:
        metricas.append(("📊 Cantidad de diapositivas", 5, "Aceptable",
                         f"{len(diaps)} de {num_esperado} esperadas"))
        total_puntos += 5
    else:
        metricas.append(("📊 Cantidad de diapositivas", 0, "Crítico", "Sin diapositivas"))

    # Calidad del contenido (30 pts más)
    max_puntos += 30
    diaps_buenas = 0
    for d in diaps:
        puntos = d.get("puntos", [])
        if d.get("titulo") and 3 <= len(puntos) <= 5:
            diaps_buenas += 1
    if diaps and diaps_buenas == len(diaps):
        metricas.append(("✨ Contenido de diapositivas", 30, "Excelente",
                         f"Todas con 3-5 puntos clave"))
        total_puntos += 30
    elif diaps_buenas > 0:
        pts = int((diaps_buenas / len(diaps)) * 30)
        metricas.append(("✨ Contenido de diapositivas", pts, "Bueno",
                         f"{diaps_buenas}/{len(diaps)} bien estructuradas"))
        total_puntos += pts
    else:
        metricas.append(("✨ Contenido de diapositivas", 0, "Deficiente",
                         "Estructura irregular"))

    # 5. Notas del orador (15 pts)
    max_puntos += 15
    con_notas = sum(1 for d in diaps if d.get("notas") and len(d["notas"]) > 20)
    if diaps and con_notas == len(diaps):
        metricas.append(("🎙️ Notas del orador", 15, "Excelente", "Todas las diapositivas"))
        total_puntos += 15
    elif con_notas > 0:
        pts = int((con_notas / len(diaps)) * 15)
        metricas.append(("🎙️ Notas del orador", pts, "Parcial", f"{con_notas}/{len(diaps)}"))
        total_puntos += pts
    else:
        metricas.append(("🎙️ Notas del orador", 0, "Ausente", "Sin notas de apoyo"))

    # 6. Conclusión y cierre (15 pts)
    max_puntos += 15
    conclusion = datos.get("conclusion", "")
    cierre = datos.get("cierre", "")
    pts_c = 0
    if conclusion and len(conclusion) >= 30: pts_c += 8
    if cierre and len(cierre) >= 10: pts_c += 7
    metricas.append(("🏁 Conclusión y cierre", pts_c,
                     "Excelente" if pts_c >= 12 else "Aceptable" if pts_c >= 6 else "Débil",
                     "Sólidos" if pts_c >= 12 else "Mejorables"))
    total_puntos += pts_c

    porcentaje = round((total_puntos / max_puntos) * 100) if max_puntos else 0
    if porcentaje >= 85:
        veredicto = "🟢 Excelente"
        color = "#10B981"
    elif porcentaje >= 65:
        veredicto = "🟡 Aceptable"
        color = "#F59E0B"
    else:
        veredicto = "🔴 Requiere mejoras"
        color = "#EF4444"

    return {
        "metricas": metricas,
        "porcentaje": porcentaje,
        "veredicto": veredicto,
        "color": color,
        "total": total_puntos,
        "max": max_puntos,
    }


def inspeccionar_infografia(datos: dict) -> dict:
    """Evalúa la calidad de la infografía generada."""
    metricas = []
    total_puntos = 0
    max_puntos = 0

    # Título + subtítulo (15 pts)
    max_puntos += 15
    titulo = datos.get("titulo", "")
    sub = datos.get("subtitulo", "")
    pts = 0
    if titulo and len(titulo) >= 8: pts += 10
    if sub and len(sub) >= 10: pts += 5
    metricas.append(("📌 Título y subtítulo", pts,
                     "Excelente" if pts >= 13 else "Aceptable" if pts >= 7 else "Débil", ""))
    total_puntos += pts

    # Introducción (10 pts)
    max_puntos += 10
    intro = datos.get("introduccion", "")
    if intro and 40 <= len(intro) <= 250:
        metricas.append(("📖 Introducción", 10, "Excelente", "Engancha al lector"))
        total_puntos += 10
    elif intro:
        metricas.append(("📖 Introducción", 5, "Aceptable", "Muy breve o extensa"))
        total_puntos += 5
    else:
        metricas.append(("📖 Introducción", 0, "Ausente", ""))

    # Estadísticas (20 pts)
    max_puntos += 20
    stats = datos.get("estadisticas", [])
    if 3 <= len(stats) <= 5:
        metricas.append(("📊 Estadísticas", 20, "Excelente", f"{len(stats)} datos destacados"))
        total_puntos += 20
    elif len(stats) > 0:
        metricas.append(("📊 Estadísticas", 10, "Parcial", f"{len(stats)} (ideal: 3-5)"))
        total_puntos += 10
    else:
        metricas.append(("📊 Estadísticas", 0, "Ausente", "Sin datos cuantitativos"))

    # Secciones (30 pts)
    max_puntos += 30
    secciones = datos.get("secciones", [])
    if 4 <= len(secciones) <= 6:
        metricas.append(("🗂️ Secciones", 30, "Excelente", f"{len(secciones)} bloques visuales"))
        total_puntos += 30
    elif len(secciones) > 0:
        metricas.append(("🗂️ Secciones", 15, "Parcial", f"{len(secciones)} (ideal: 4-6)"))
        total_puntos += 15
    else:
        metricas.append(("🗂️ Secciones", 0, "Ausente", ""))

    # Conclusión (15 pts)
    max_puntos += 15
    concl = datos.get("conclusion", "")
    if concl and len(concl) >= 30:
        metricas.append(("🏁 Conclusión", 15, "Excelente", "Cierre con impacto"))
        total_puntos += 15
    elif concl:
        metricas.append(("🏁 Conclusión", 7, "Aceptable", "Breve"))
        total_puntos += 7
    else:
        metricas.append(("🏁 Conclusión", 0, "Ausente", ""))

    # Timeline opcional (10 pts bonus)
    max_puntos += 10
    timeline = datos.get("timeline", [])
    if len(timeline) >= 3:
        metricas.append(("⏱️ Línea de tiempo", 10, "Bonus", f"{len(timeline)} etapas"))
        total_puntos += 10
    elif len(timeline) > 0:
        metricas.append(("⏱️ Línea de tiempo", 5, "Parcial", f"{len(timeline)} etapas"))
        total_puntos += 5
    else:
        metricas.append(("⏱️ Línea de tiempo", 0, "Opcional", "No aplica al tema"))

    porcentaje = round((total_puntos / max_puntos) * 100) if max_puntos else 0
    if porcentaje >= 85:
        veredicto = "🟢 Excelente"
        color = "#10B981"
    elif porcentaje >= 65:
        veredicto = "🟡 Aceptable"
        color = "#F59E0B"
    else:
        veredicto = "🔴 Requiere mejoras"
        color = "#EF4444"

    return {
        "metricas": metricas,
        "porcentaje": porcentaje,
        "veredicto": veredicto,
        "color": color,
        "total": total_puntos,
        "max": max_puntos,
    }

# ═══════════════════════════════════════════════════════════════════════════
# UTILIDADES
# ═══════════════════════════════════════════════════════════════════════════
def extraer_texto_pdf(archivo, max_caracteres=60000):
    if PdfReader is None:
        raise RuntimeError("No hay librería PDF disponible. Instala pypdf o PyPDF2.")
    archivo.seek(0)
    reader = PdfReader(archivo)
    texto = "".join((p.extract_text() or "") for p in reader.pages)
    texto = re.sub(r"\s+", " ", texto).strip()
    return texto[:max_caracteres]


def construir_contexto(archivo_pdf, texto_base):
    contenido = ""
    if archivo_pdf is not None:
        try:
            contenido = extraer_texto_pdf(archivo_pdf)
        except Exception:
            contenido = ""
    if not contenido and texto_base:
        contenido = texto_base.strip()
    if contenido:
        return ("CONTENIDO CURRICULAR DE REFERENCIA (base única, sé fiel a él):\n"
                f"{contenido}\n")
    return "SIN CONTENIDO DE REFERENCIA: genera contenido riguroso y preciso sobre el tema.\n"

# ═══════════════════════════════════════════════════════════════════════════
# CONSTRUCTOR DE POWERPOINT (con imágenes IA)
# ═══════════════════════════════════════════════════════════════════════════
def build_pptx(datos, meta, tema_key, imagenes=None):
    tema = TEMAS_PPT.get(tema_key, TEMAS_PPT["🔵 Institucional Azul"])
    SW, SH = 13.333, 7.5
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    blank = prs.slide_layouts[6]

    def hexc(h):
        return PColor.from_string(h)

    def add_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = hexc(color)
        bg.line.fill.background()
        bg.shadow.inherit = False

    def add_rect(slide, l, t, w, h, color):
        r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        r.fill.solid()
        r.fill.fore_color.rgb = hexc(color)
        r.line.fill.background()
        r.shadow.inherit = False
        return r

    def add_text(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        p.text = text
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = hexc(color)
            run.font.name = "Calibri"
        return box

    def add_bullets(slide, l, t, w, h, points, size, color):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, pt in enumerate(points):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"▪  {pt}"
            p.space_after = Pt(10)
            for run in p.runs:
                run.font.size = Pt(size)
                run.font.color.rgb = hexc(color)
                run.font.name = "Calibri"
        return box

    # ── Diapositiva 1: PORTADA con imagen de fondo ──
    s = prs.slides.add_slide(blank)
    add_bg(s, tema["bg"])
    if imagenes and imagenes.get("portada"):
        try:
            s.shapes.add_picture(imagenes["portada"], Inches(6.5), Inches(0.8),
                                 width=Inches(6.3), height=Inches(3.6))
        except Exception:
            pass
    add_rect(s, 0, 0, 0.35, SH, tema["accent"])
    add_rect(s, 0.7, 0.8, 5.8, 5.8, tema["bg"])
    add_text(s, 1.0, 1.2, 5.3, 0.5, meta.get("asignatura", "").upper(), 16, tema["accent"], bold=True)
    add_text(s, 1.0, 2.0, 5.3, 2.2, datos.get("titulo", "Presentación"), 40, tema["text"], bold=True)
    add_text(s, 1.0, 4.4, 5.3, 1.0, datos.get("subtitulo", ""), 18, tema["sub"])
    add_rect(s, 1.0, 5.6, 3.0, 0.06, tema["accent"])
    add_text(s, 1.0, 5.9, 5.3, 0.5,
             f'{meta.get("docente", "")}  ·  {meta.get("centro", "")}', 13, tema["sub"])

    # ── Diapositiva 2: Objetivo + Agenda ──
    s = prs.slides.add_slide(blank)
    add_bg(s, tema["bg"])
    add_rect(s, 0, 0, SW, 0.25, tema["accent"])
    add_text(s, 0.8, 0.7, 11, 0.8, "🎯 Objetivo de Aprendizaje", 30, tema["text"], bold=True)
    add_rect(s, 0.8, 1.6, 2.2, 0.05, tema["accent"])
    add_text(s, 0.8, 2.0, 11.7, 1.6, datos.get("objetivo", ""), 20, tema["sub"])
    add_text(s, 0.8, 4.0, 11, 0.6, "Contenido de la sesión", 18, tema["accent"], bold=True)
    agenda = [d.get("titulo", "") for d in datos.get("diapositivas", [])][:6]
    add_bullets(s, 0.8, 4.6, 11.7, 2.4, agenda, 15, tema["text"])

    # ── Diapositivas de CONTENIDO con imagen lateral ──
    for i, diap in enumerate(datos.get("diapositivas", []), 1):
        s = prs.slides.add_slide(blank)
        add_bg(s, tema["bg"])
        add_rect(s, 0, 0, SW, 0.25, tema["accent"])
        add_rect(s, 0.8, 0.7, 0.9, 0.9, tema["accent"])
        add_text(s, 0.8, 0.82, 0.9, 0.7, str(i), 28, "FFFFFF", bold=True, align=PP_ALIGN.CENTER)
        add_text(s, 2.0, 0.85, 6.5, 0.8, diap.get("titulo", ""), 26, tema["text"], bold=True)
        add_rect(s, 0.8, 1.9, 7.0, 0.03, tema["accent"])
        add_bullets(s, 0.9, 2.3, 6.8, 4.2, diap.get("puntos", []), 16, tema["text"])
        # Imagen a la derecha
        if imagenes and imagenes.get(f"diap_{i}"):
            try:
                s.shapes.add_picture(imagenes[f"diap_{i}"],
                                     Inches(8.3), Inches(1.8),
                                     width=Inches(4.5), height=Inches(4.5))
            except Exception:
                pass
        if diap.get("notas"):
            s.notes_slide.notes_text_frame.text = diap.get("notas", "")

    # ── Conclusión con imagen ──
    s = prs.slides.add_slide(blank)
    add_bg(s, tema["bg"])
    add_rect(s, 0, 0, SW, 0.25, tema["accent"])
    if imagenes and imagenes.get("conclusion"):
        try:
            s.shapes.add_picture(imagenes["conclusion"],
                                 Inches(7.5), Inches(1.5),
                                 width=Inches(5.3), height=Inches(4.5))
        except Exception:
            pass
    add_text(s, 0.8, 0.8, 6.5, 0.8, "✅ Conclusión", 30, tema["text"], bold=True)
    add_rect(s, 0.8, 1.7, 2.2, 0.05, tema["accent"])
    add_text(s, 0.8, 2.1, 6.5, 4.5, datos.get("conclusion", ""), 18, tema["sub"])

    # ── Cierre ──
    s = prs.slides.add_slide(blank)
    add_bg(s, tema["bg"])
    add_rect(s, 0, 0, 0.35, SH, tema["accent"])
    add_text(s, 1.0, 2.4, 11.3, 1.5, "¡Gracias!", 54, tema["text"],
             bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 1.0, 4.2, 11.3, 1.0,
             datos.get("cierre", "¿Preguntas o comentarios?"), 20, tema["sub"],
             align=PP_ALIGN.CENTER)
    add_text(s, 1.0, 5.6, 11.3, 0.5,
             f'{meta.get("docente", "")} · {meta.get("centro", "")}', 14, tema["sub"],
             align=PP_ALIGN.CENTER)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ═══════════════════════════════════════════════════════════════════════════
# CONSTRUCTOR DE INFOGRAFÍA HTML (con imágenes IA)
# ═══════════════════════════════════════════════════════════════════════════
def build_infografia_html(datos, meta, tema_key, imagenes=None):
    tema = TEMAS_INFO.get(tema_key, TEMAS_INFO["🔵 Azul Institucional"])
    esc = _html.escape

    # Hero con imagen
    hero_img_html = ""
    if imagenes and imagenes.get("hero"):
        hero_img_html = (f'<img src="{esc(imagenes["hero"])}" '
                         f'style="width:100%; max-height:340px; object-fit:cover; '
                         f'border-radius:20px; margin-bottom:20px; '
                         f'box-shadow:0 10px 30px rgba(0,0,0,0.15);" '
                         f'alt="Imagen del tema">')

    secciones = ""
    for i, sec in enumerate(datos.get("secciones", []), 1):
        img_html = ""
        if imagenes and imagenes.get(f"sec_{i}"):
            img_html = (f'<img src="{esc(imagenes[f"sec_{i}"])}" '
                        f'style="width:100%; height:160px; object-fit:cover; '
                        f'border-radius:12px; margin-bottom:12px;" alt="">')
        secciones += (
            '<div class="card">'
            f'{img_html}'
            f'<div class="card-icon">{esc(str(sec.get("icono", "✨")))}</div>'
            f'<div class="card-title">{esc(str(sec.get("titulo", "🚀")))}</div>'
            f'<div class="card-desc">{esc(str(sec.get("descripcion", "")))}</div>'
            '</div>'
        )

    stats = ""
    for stt in datos.get("estadisticas", []):
        stats += (
            '<div class="stat">'
            f'<div class="stat-value">{esc(str(stt.get("valor", "")))}</div>'
            f'<div class="stat-label">{esc(str(stt.get("etiqueta", "")))}</div>'
            '</div>'
        )

    timeline = ""
    for tl in datos.get("timeline", []):
        timeline += (
            '<div class="tl-item">'
            '<div class="tl-dot"></div>'
            '<div class="tl-body">'
            f'<div class="tl-stage">{esc(str(tl.get("etapa", "")))}</div>'
            f'<div class="tl-desc">{esc(str(tl.get("descripcion", "")))}</div>'
            '</div></div>'
        )

    css_template = """
* { margin:0; padding:0; box-sizing:border-box; }
body { font-family:'Segoe UI', Arial, sans-serif; background: BG ; color: TEXT ; padding:40px 20px; }
.wrap { max-width:900px; margin:0 auto; }
.header { text-align:center; padding:44px 24px; background:linear-gradient(135deg, ACCENT , ACCENT2 );
          border-radius:20px; color:#fff; margin-bottom:28px;
          box-shadow:0 10px 30px rgba(0,0,0,0.15); }
.header h1 { font-size:36px; margin-bottom:10px; line-height:1.2; }
.header p { font-size:18px; opacity:0.95; }
.intro { background: CARD ; border-radius:16px; padding:24px; margin-bottom:26px;
         box-shadow:0 4px 16px rgba(0,0,0,0.06); font-size:16px; line-height:1.6;
         border-left:6px solid ACCENT ; }
.stats { display:flex; gap:16px; flex-wrap:wrap; margin-bottom:26px; }
.stat { flex:1; min-width:150px; background: CARD ; border-radius:16px; padding:24px;
        text-align:center; box-shadow:0 4px 16px rgba(0,0,0,0.06);
        border-top:5px solid ACCENT ; }
.stat-value { font-size:34px; font-weight:800; color: ACCENT ; }
.stat-label { font-size:14px; color: MUTED ; margin-top:6px; }
.grid { display:grid; grid-template-columns:repeat(auto-fit,minmax(250px,1fr));
        gap:18px; margin-bottom:26px; }
.card { background: CARD ; border-radius:16px; padding:24px;
        box-shadow:0 4px 16px rgba(0,0,0,0.06); border-top:5px solid ACCENT2 ;
        transition: transform 0.25s ease; }
.card:hover { transform: translateY(-4px); box-shadow:0 12px 28px rgba(0,0,0,0.12); }
.card-icon { font-size:32px; margin-bottom:10px; }
.card-title { font-size:18px; font-weight:700; margin-bottom:8px; color: TEXT ; }
.card-desc { font-size:14px; line-height:1.5; color: MUTED ; }
.timeline { margin-bottom:26px; background: CARD ; border-radius:16px; padding:24px;
            box-shadow:0 4px 16px rgba(0,0,0,0.06); }
.tl-item { display:flex; gap:16px; margin-bottom:18px; }
.tl-item:last-child { margin-bottom:0; }
.tl-dot { width:16px; height:16px; border-radius:50%; background: ACCENT ;
          margin-top:4px; flex-shrink:0; box-shadow:0 0 0 5px ACCENTSOFT ; }
.tl-stage { font-weight:700; color: TEXT ; margin-bottom:4px; }
.tl-desc { font-size:14px; color: MUTED ; line-height:1.5; }
.footer { text-align:center; padding:26px; background: CARD ; border-radius:16px;
          box-shadow:0 4px 16px rgba(0,0,0,0.06); }
.footer .concl { font-size:16px; line-height:1.6; margin-bottom:14px; }
.footer .cred { font-size:13px; color: MUTED ; }
"""
    css = (css_template
           .replace("ACCENT2", tema["accent2"])
           .replace("ACCENTSOFT", tema["accentsoft"])
           .replace("ACCENT", tema["accent"])
           .replace("BG", tema["bg"])
           .replace("CARD", tema["card"])
           .replace("MUTED", tema["muted"])
           .replace("TEXT", tema["text"]))

    html_doc = f"""<!DOCTYPE html>
<html lang="es">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{esc(str(datos.get('titulo', 'Infografía')))}</title>
    <style>{css}</style>
</head>
<body>
<div class="wrap">
    {hero_img_html}
    <div class="header">
        <h1>{esc(str(datos.get('titulo', '')))}</h1>
        <p>{esc(str(datos.get('subtitulo', '')))}</p>
    </div>
    <div class="intro">{esc(str(datos.get('introduccion', '')))}</div>
    {('<div class="stats">' + stats + '</div>') if stats else ''}
    {('<div class="grid">' + secciones + '</div>') if secciones else ''}
    {('<div class="timeline">' + timeline + '</div>') if timeline else ''}
    <div class="footer">
        <div class="concl">{esc(str(datos.get('conclusion', '')))}</div>
        <div class="cred">{esc(str(meta.get('docente', '')))} · {esc(str(meta.get('centro', '')))} · {esc(str(meta.get('asignatura', '')))}</div>
    </div>
</div>
</body>
</html>"""
    return html_doc

# ═══════════════════════════════════════════════════════════════════════════
# PROMPTS
# ═══════════════════════════════════════════════════════════════════════════
def prompt_ppt(form, contexto):
    return f"""Actúa como un Diseñador Instruccional experto en presentaciones educativas de la ETP del MINERD.
Genera el contenido para una presentación de diapositivas profesional.
DATOS:
Tema: {form['tema']}
Asignatura/Módulo: {form['asignatura']}
Nivel educativo: {form['nivel']}
Audiencia: {form['audiencia']}
Tono: {form['tono']}
Cantidad EXACTA de diapositivas de contenido: {form['num_diapositivas']}
{contexto}
REGLAS:
• Cada diapositiva: título corto (máx 8 palabras) y entre 3 y 5 puntos clave (máx. 20 palabras cada una).
• Notas del orador por diapositiva (2-3 oraciones de apoyo).
• Objetivo claro, medible, con verbo de la taxonomía de Bloom.
• Conclusión sintética (máx. 60 palabras) y frase de cierre motivadora.
CODIFICACIÓN OBLIGATORIA: salto de línea → {ia.MARKER_NL} · comilla doble → {ia.MARKER_DQ} · tabulación → {ia.MARKER_TAB}.
Devuelve SOLO JSON válido:
{{
 "titulo": "Título principal",
 "subtitulo": "Subtítulo atractivo",
 "objetivo": "Objetivo de aprendizaje con verbo observable",
 "diapositivas": [
    {{"titulo": "Título", "puntos": ["P1", "P2", "P3"], "notas": "Lo que dirá el docente"}}
 ],
 "conclusion": "Conclusión que sintetiza el aprendizaje",
 "cierre": "Frase de cierre o invitación a preguntas"
}}
El arreglo "diapositivas" debe tener EXACTAMENTE {form['num_diapositivas']} elementos."""


def prompt_infografia(form, contexto):
    return f"""Actúa como un Diseñador de Infografías Educativas experto en síntesis visual para la ETP del MINERD.
Genera el contenido estructurado para una infografía impactante y clara.
DATOS:
Tema: {form['tema']}
Asignatura/Módulo: {form['asignatura']}
Nivel educativo: {form['nivel']}
Audiencia: {form['audiencia']}
Tono: {form['tono']}
{contexto}
REGLAS:
• Contenido sintético, visual y fácil de escanear.
• 3-4 estadísticas destacadas (valor corto y etiqueta breve).
• 4-6 secciones con emoji, título y descripción corta (máx. 30 palabras).
• Opcionalmente timeline de 3-5 etapas si el tema lo amerita.
CODIFICACIÓN OBLIGATORIA: salto de línea → {ia.MARKER_NL} · comilla doble → {ia.MARKER_DQ} · tabulación → {ia.MARKER_TAB}.
Devuelve SOLO JSON válido:
{{
 "titulo": "Título impactante",
 "subtitulo": "Subtítulo breve",
 "introduccion": "Introducción de 2 a 3 líneas que engancha",
 "estadisticas": [{{"valor": "85%", "etiqueta": "Descripción"}}],
 "secciones": [{{"icono": "🎯", "titulo": "Título", "descripcion": "Descripción"}}],
 "timeline": [{{"etapa": "Etapa 1", "descripcion": "Descripción"}}],
 "conclusion": "Mensaje de cierre"
}}"""

# ═══════════════════════════════════════════════════════════════════════════
# ESTILOS NIVEL DIOS
# ═══════════════════════════════════════════════════════════════════════════
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800;900&display=swap');
html, body, [class*="css"] { font-family: 'Inter', sans-serif; background-color: #F8FAFC; color: #1E293B; }
.visual-hero {
  background: linear-gradient(135deg, #0F172A 0%, #7C3AED 55%, #EC4899 100%);
  color:#fff; padding:2.4rem; border-radius:20px; margin-bottom:1.5rem;
  box-shadow:0 25px 50px rgba(124,58,237,0.35); position:relative; overflow:hidden;
}
.visual-hero::before {
  content:''; position:absolute; top:-50%; left:-50%; width:200%; height:200%;
  background:radial-gradient(circle, rgba(255,255,255,0.15) 0%, transparent 60%);
  animation:vpulse 6s ease-in-out infinite;
}
@keyframes vpulse { 0%,100%{transform:scale(1);opacity:.5} 50%{transform:scale(1.1);opacity:.9} }
.visual-hero-title { font-size:2.5rem; font-weight:900; letter-spacing:-0.03em; position:relative; }
.visual-hero-sub { font-size:1.08rem; opacity:0.92; margin-top:0.4rem; position:relative; line-height:1.5; }
.visual-hero-badge { display:inline-block; background:rgba(255,255,255,0.18);
  border:1px solid rgba(255,255,255,0.3); border-radius:8px; padding:5px 14px;
  font-size:0.82rem; font-weight:600; margin-top:0.9rem; margin-right:8px; position:relative; }
.visual-section-title { color:#7C3AED; font-weight:700; font-size:1.15rem;
  border-bottom:2px solid #EDE9FE; padding-bottom:8px; margin:1.3rem 0 1rem 0; }

.inspection-card {
  background: linear-gradient(135deg, #FFFFFF 0%, #F8FAFC 100%);
  border:2px solid #E2E8F0; border-radius:16px; padding:1.4rem;
  box-shadow:0 8px 24px rgba(0,0,0,0.08); margin-bottom:1rem;
}
.inspection-header { display:flex; justify-content:space-between; align-items:center;
  margin-bottom:1rem; padding-bottom:0.8rem; border-bottom:2px dashed #E2E8F0; }
.inspection-score { font-size:3rem; font-weight:900; line-height:1; }
.inspection-badge { display:inline-block; padding:6px 14px; border-radius:999px;
  font-size:0.85rem; font-weight:700; letter-spacing:0.02em; }
.metric-row { display:flex; justify-content:space-between; align-items:center;
  padding:8px 12px; border-radius:8px; margin-bottom:4px; transition: background 0.2s; }
.metric-row:hover { background:#F1F5F9; }
.metric-name { font-size:0.9rem; font-weight:600; color:#1E293B; flex:1; }
.metric-status { font-size:0.78rem; font-weight:700; padding:3px 10px;
  border-radius:6px; margin-left:8px; }
.metric-detail { font-size:0.78rem; color:#64748B; margin-left:8px; }

.image-gallery { display:grid; grid-template-columns:repeat(auto-fill,minmax(220px,1fr));
  gap:14px; margin:1rem 0; }
.image-tile { position:relative; border-radius:12px; overflow:hidden;
  box-shadow:0 4px 12px rgba(0,0,0,0.1); background:#fff; transition: all 0.3s; }
.image-tile:hover { transform:translateY(-4px); box-shadow:0 12px 28px rgba(0,0,0,0.18); }
.image-tile img { width:100%; height:150px; object-fit:cover; display:block; }
.image-tile-label { padding:8px 12px; font-size:0.82rem; font-weight:600;
  color:#1E293B; background:#fff; border-top:1px solid #E2E8F0; }
.image-tile-overlay { position:absolute; top:8px; right:8px; background:rgba(15,23,42,0.85);
  color:#fff; font-size:0.7rem; padding:3px 8px; border-radius:6px; font-weight:600; }
</style>
""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════════════════
# GUARDIA + ENCABEZADO
# ═══════════════════════════════════════════════════════════════════════════
if not (st.session_state.get("docente_autenticado") or st.session_state.get("coordinador_autenticado")):
    st.error("🔒 Inicia sesión para crear presentaciones e infografías.")
    st.stop()

ia.panel_sidebar_ia("Fábrica Visual IA")
st.markdown("""
<div class="visual-hero">
    <div class="visual-hero-title">🎨 Fábrica Visual IA · Nivel Dios</div>
    <div class="visual-hero-sub">
        Presentaciones PowerPoint e infografías con <b>imágenes IA automáticas</b> y
        <b>auto-inspección inteligente</b>. Cada producto se audita a sí mismo antes de entregarse.
    </div>
    <div>
        <span class="visual-hero-badge">🖼️ Imágenes IA (Pollinations)</span>
        <span class="visual-hero-badge">🔍 Auto-Inspección</span>
        <span class="visual-hero-badge">📊 PowerPoint 16:9</span>
        <span class="visual-hero-badge">🎨 Infografías HTML</span>
        <span class="visual-hero-badge">🔄 Galería Regenerable</span>
    </div>
</div>
""", unsafe_allow_html=True)

if not PPTX_OK:
    st.warning("⚠️ Para generar presentaciones PowerPoint instala: `pip install python-pptx`")

# Estado persistente
st.session_state.setdefault("ppt_result", None)
st.session_state.setdefault("inf_result", None)
st.session_state.setdefault("ppt_images", None)
st.session_state.setdefault("inf_images", None)
st.session_state.setdefault("ppt_inspection", None)
st.session_state.setdefault("inf_inspection", None)

nombre_docente = st.session_state.get("nombre_docente") or "Ing. Bernardo Antonio Hernández Batista"
tab_ppt, tab_inf = st.tabs(["📊 Presentación PowerPoint", "🎨 Infografía"])

# ═══════════════════════════════════════════════════════════════════════════
# RENDER PANEL DE INSPECCIÓN
# ═══════════════════════════════════════════════════════════════════════════
def render_panel_inspeccion(insp: dict):
    color = insp["color"]
    st.markdown(f"""
    <div class="inspection-card">
        <div class="inspection-header">
            <div>
                <div style="font-size:0.85rem; font-weight:700; color:#64748B;
                            text-transform:uppercase; letter-spacing:0.05em;">
                    🔍 Auto-Inspección IA
                </div>
                <div style="font-size:1.4rem; font-weight:800; color:#0F172A; margin-top:4px;">
                    {insp['veredicto']}
                </div>
            </div>
            <div style="text-align:right;">
                <div class="inspection-score" style="color:{color};">
                    {insp['porcentaje']}<span style="font-size:1.2rem;">%</span>
                </div>
                <div style="font-size:0.75rem; color:#94A3B8; font-weight:600;">
                    {insp['total']}/{insp['max']} pts
                </div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    for nombre, pts, estado, detalle in insp["metricas"]:
        color_est = {"Excelente": "#10B981", "Bueno": "#10B981",
                     "Aceptable": "#F59E0B", "Parcial": "#F59E0B",
                     "Débil": "#EF4444", "Ausente": "#EF4444", "Crítico": "#EF4444",
                     "Deficiente": "#EF4444", "Bonus": "#8B5CF6", "Opcional": "#94A3B8"
                    }.get(estado, "#64748B")
        bg_est = {"Excelente": "#D1FAE5", "Bueno": "#D1FAE5",
                  "Aceptable": "#FEF3C7", "Parcial": "#FEF3C7",
                  "Débil": "#FEE2E2", "Ausente": "#FEE2E2", "Crítico": "#FEE2E2",
                  "Deficiente": "#FEE2E2", "Bonus": "#EDE9FE", "Opcional": "#F1F5F9"
                 }.get(estado, "#F1F5F9")
        st.markdown(f"""
        <div class="metric-row">
            <span class="metric-name">{nombre}</span>
            <span class="metric-detail">{detalle}</span>
            <span class="metric-status" style="background:{bg_est}; color:{color_est};">
                {estado} · {pts} pts
            </span>
        </div>
        """, unsafe_allow_html=True)


def render_galeria_imagenes(imagenes: dict, pref: str):
    """Muestra la galería con botón de regenerar individual."""
    st.markdown('<div class="visual-section-title">🖼️ Galería de Imágenes IA</div>',
                unsafe_allow_html=True)
    st.caption("💡 Cada imagen es única y generada por IA. Puedes regenerar cualquiera con un clic.")
    cols = st.columns(min(4, len(imagenes)))
    keys = list(imagenes.keys())
    for idx, key_img in enumerate(keys):
        col = cols[idx % len(cols)]
        with col:
            st.image(imagenes[key_img], use_container_width=True)
            etiqueta = {"portada": "🎬 Portada", "conclusion": "🏁 Conclusión",
                        "hero": "🎨 Hero", "cierre": "👋 Cierre"}.get(key_img,
                        f"📄 {key_img.replace('_', ' ').title()}")
            st.caption(f"**{etiqueta}**")
            if st.button("🔄 Regenerar", key=f"regen_{pref}_{key_img}",
                         use_container_width=True):
                new_seed = random.randint(1, 99999)
                if pref == "ppt":
                    datos = st.session_state.ppt_result["datos"]
                    tema = st.session_state.ppt_result.get("tema_diseno", "🔵 Institucional Azul")
                    if key_img == "portada":
                        new_url = generar_url_imagen(
                            f"Abstract background for {datos.get('titulo','')}",
                            1280, 720, seed=new_seed, estilo="minimalist")
                    elif key_img == "conclusion":
                        new_url = generar_url_imagen(
                            f"Closing visual for {datos.get('titulo','')}",
                            1024, 576, seed=new_seed, estilo="3d")
                    elif key_img.startswith("diap_"):
                        i = int(key_img.split("_")[1])
                        d = datos.get("diapositivas", [])[i-1] if i <= len(datos.get("diapositivas", [])) else {}
                        new_url = generar_url_imagen(
                            f"Educational visual for {d.get('titulo','')}",
                            800, 500, seed=new_seed, estilo="academic")
                    else:
                        new_url = imagenes[key_img]
                    st.session_state.ppt_images[key_img] = new_url
                else:
                    datos = st.session_state.inf_result["datos"]
                    if key_img == "hero":
                        new_url = generar_url_imagen(
                            f"Hero banner for {datos.get('titulo','')}",
                            1200, 400, seed=new_seed, estilo="illustration")
                    elif key_img.startswith("sec_"):
                        i = int(key_img.split("_")[1])
                        s = datos.get("secciones", [])[i-1] if i <= len(datos.get("secciones", [])) else {}
                        new_url = generar_url_imagen(
                            f"Illustration for {s.get('titulo','')}",
                            400, 400, seed=new_seed, estilo="illustration")
                    else:
                        new_url = imagenes[key_img]
                    st.session_state.inf_images[key_img] = new_url
                st.toast(f"✅ Imagen '{key_img}' regenerada", icon="🔄")
                st.rerun()

# ═══════════════════════════════════════════════════════════════════════════
# TAB 1: POWERPOINT
# ═══════════════════════════════════════════════════════════════════════════
with tab_ppt:
    with st.form("form_ppt", clear_on_submit=False):
        st.markdown('<div class="visual-section-title">📝 1. Contenido</div>', unsafe_allow_html=True)
        tema_ppt = st.text_input("Tema de la presentación *",
                                  placeholder="Ej: Simetría central con respecto al origen")
        col1, col2 = st.columns(2)
        with col1:
            asignatura_ppt = st.text_input("Asignatura / Módulo", placeholder="Ej: Geometría / MF 358-3")
            nivel_ppt = st.selectbox("Nivel educativo",
                ["Secundaria", "Bachillerato Técnico (ETP)", "Primaria", "Universitario"])
        with col2:
            audiencia_ppt = st.text_input("Audiencia", placeholder="Ej: Estudiantes de 4to B")
            tono_ppt = st.selectbox("Tono", TONOS)
        col3, col4 = st.columns(2)
        with col3:
            num_diapositivas = st.slider("Diapositivas de contenido", 3, 12, 6)
        with col4:
            tema_diseno_ppt = st.selectbox("Tema de diseño", list(TEMAS_PPT.keys()))

        st.markdown('<div class="visual-section-title">📎 2. Anclaje curricular (opcional)</div>',
                    unsafe_allow_html=True)
        archivo_ppt = st.file_uploader("PDF de referencia", type=["pdf"], key="pdf_ppt")
        texto_base_ppt = st.text_area("O pega el contenido base", height=80,
            placeholder="Pega aquí el contenido curricular o los puntos clave...", key="txt_ppt")

        chk_img_ppt = st.checkbox("🖼️ Generar imágenes IA automáticas (recomendado)", value=True, key="chk_img_ppt")
        chk_insp_ppt = st.checkbox("🔍 Auto-inspección inteligente (recomendado)", value=True, key="chk_insp_ppt")

        max_tokens, temperature = ia.control_avanzado(default_tokens=12000, tope=32000, default_temp=0.4)
        submit_ppt = st.form_submit_button("✨ Generar Presentación + Imágenes IA",
                                           type="primary", width="stretch")

    if submit_ppt:
        if not tema_ppt.strip():
            st.warning("📝 Escribe el tema de la presentación.")
        elif not PPTX_OK:
            st.error("⚠️ Instala python-pptx: `pip install python-pptx`")
        else:
            with st.spinner("🧠 Diseñando presentación con IA..."):
                try:
                    contexto = construir_contexto(archivo_ppt, texto_base_ppt)
                    form = {
                        "tema": tema_ppt, "asignatura": asignatura_ppt, "nivel": nivel_ppt,
                        "audiencia": audiencia_ppt, "tono": tono_ppt,
                        "num_diapositivas": num_diapositivas,
                    }
                    prompt = prompt_ppt(form, contexto)
                    texto_crudo, flags = ia.solicitar_ia(
                        prompt, modo="json", max_tokens=max_tokens,
                        temperature=temperature, modulo="presentaciones_ppt")
                    datos = ia.decodificar_marcadores(ia.parsear_json_robusto(texto_crudo))

                    # Generar imágenes IA
                    imagenes = None
                    if chk_img_ppt:
                        with st.spinner("🖼️ Generando imágenes IA por diapositiva..."):
                            imagenes = generar_imagenes_para_presentacion(datos, tema_diseno_ppt)
                        st.session_state.ppt_images = imagenes

                    # Auto-inspección
                    inspeccion = None
                    if chk_insp_ppt:
                        inspeccion = inspeccionar_presentacion(datos, num_diapositivas)
                        st.session_state.ppt_inspection = inspeccion

                    meta = {"docente": nombre_docente,
                            "centro": "Politécnico Salesiano Arquides Calderón",
                            "asignatura": asignatura_ppt}
                    buffer = build_pptx(datos, meta, tema_diseno_ppt, imagenes)
                    st.session_state.ppt_result = {
                        "buffer": buffer, "datos": datos, "tema": tema_ppt,
                        "tema_diseno": tema_diseno_ppt, "meta": meta,
                    }
                    st.toast("✅ Presentación generada con imágenes IA.", icon="🎨")
                    st.rerun()
                except Exception as e:
                    ia.render_error_ia(e, None)

    # RENDER RESULTADO PPT
    if st.session_state.ppt_result:
        res = st.session_state.ppt_result
        datos = res["datos"]

        # Auto-inspección
        if st.session_state.ppt_inspection:
            st.markdown('<div class="visual-section-title">🔍 Auto-Inspección</div>',
                        unsafe_allow_html=True)
            render_panel_inspeccion(st.session_state.ppt_inspection)

        # Galería de imágenes
        if st.session_state.ppt_images:
            render_galeria_imagenes(st.session_state.ppt_images, "ppt")

        # Estructura
        st.markdown('<div class="visual-section-title">📋 Estructura generada</div>',
                    unsafe_allow_html=True)
        c1, c2 = st.columns([2, 1])
        with c1:
            st.markdown(f"### {datos.get('titulo', '')}")
            st.caption(datos.get("subtitulo", ""))
            st.info(f"🎯 **Objetivo:** {datos.get('objetivo', '')}")
            with st.expander("🗂️ Ver estructura completa"):
                for i, d in enumerate(datos.get("diapositivas", []), 1):
                    st.markdown(f"**{i}. {d.get('titulo', '')}**")
                    for p in d.get("puntos", []):
                        st.markdown(f"&nbsp;&nbsp;• {p}")
                    if d.get("notas"):
                        st.caption(f"🎙️ {d['notas']}")
        with c2:
            st.download_button(
                label="📥 Descargar PowerPoint (.pptx)",
                data=res["buffer"],
                file_name=f"{ia.sanear_nombre_archivo(res['tema'])}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary", width="stretch")
            if st.button("🗑️ Descartar presentación", width="stretch", key="discard_ppt"):
                for k in ["ppt_result", "ppt_images", "ppt_inspection"]:
                    st.session_state[k] = None
                st.rerun()

# ═══════════════════════════════════════════════════════════════════════════
# TAB 2: INFOGRAFÍA
# ═══════════════════════════════════════════════════════════════════════════
with tab_inf:
    with st.form("form_inf", clear_on_submit=False):
        st.markdown('<div class="visual-section-title">📝 1. Contenido</div>', unsafe_allow_html=True)
        tema_inf = st.text_input("Tema de la infografía *",
                                  placeholder="Ej: Proceso de importación vehicular")
        col1, col2 = st.columns(2)
        with col1:
            asignatura_inf = st.text_input("Asignatura / Módulo", placeholder="Ej: MF 358-3")
            nivel_inf = st.selectbox("Nivel educativo",
                ["Secundaria", "Bachillerato Técnico (ETP)", "Primaria", "Universitario"])
        with col2:
            audiencia_inf = st.text_input("Audiencia", placeholder="Ej: Estudiantes de 5to B")
            tono_inf = st.selectbox("Tono", TONOS)
        tema_diseno_inf = st.selectbox("Tema de diseño", list(TEMAS_INFO.keys()))

        st.markdown('<div class="visual-section-title">📎 2. Anclaje curricular (opcional)</div>',
                    unsafe_allow_html=True)
        archivo_inf = st.file_uploader("PDF de referencia", type=["pdf"], key="pdf_inf")
        texto_base_inf = st.text_area("O pega el contenido base", height=80,
            placeholder="Pega aquí el contenido curricular...", key="txt_inf")

        chk_img_inf = st.checkbox("🖼️ Generar imágenes IA automáticas (recomendado)", value=True, key="chk_img_inf")
        chk_insp_inf = st.checkbox("🔍 Auto-inspección inteligente (recomendado)", value=True, key="chk_insp_inf")

        max_tokens_i, temperature_i = ia.control_avanzado(default_tokens=10000, tope=32000, default_temp=0.5)
        submit_inf = st.form_submit_button("✨ Generar Infografía + Imágenes IA",
                                           type="primary", width="stretch")

    if submit_inf:
        if not tema_inf.strip():
            st.warning("📝 Escribe el tema de la infografía.")
        else:
            with st.spinner("🧠 Diseñando infografía con IA..."):
                try:
                    contexto = construir_contexto(archivo_inf, texto_base_inf)
                    form = {
                        "tema": tema_inf, "asignatura": asignatura_inf, "nivel": nivel_inf,
                        "audiencia": audiencia_inf, "tono": tono_inf,
                    }
                    prompt = prompt_infografia(form, contexto)
                    texto_crudo, flags = ia.solicitar_ia(
                        prompt, modo="json", max_tokens=max_tokens_i,
                        temperature=temperature_i, modulo="presentaciones_infografia")
                    datos = ia.decodificar_marcadores(ia.parsear_json_robusto(texto_crudo))

                    imagenes = None
                    if chk_img_inf:
                        with st.spinner("🖼️ Generando imágenes IA para secciones..."):
                            imagenes = generar_imagenes_para_infografia(datos, tema_diseno_inf)
                        st.session_state.inf_images = imagenes

                    inspeccion = None
                    if chk_insp_inf:
                        inspeccion = inspeccionar_infografia(datos)
                        st.session_state.inf_inspection = inspeccion

                    meta = {"docente": nombre_docente,
                            "centro": "Politécnico Salesiano Arquides Calderón",
                            "asignatura": asignatura_inf}
                    html_doc = build_infografia_html(datos, meta, tema_diseno_inf, imagenes)
                    st.session_state.inf_result = {
                        "html": html_doc, "tema": tema_inf, "datos": datos, "meta": meta,
                    }
                    st.toast("✅ Infografía generada con imágenes IA.", icon="🎨")
                    st.rerun()
                except Exception as e:
                    ia.render_error_ia(e, None)

    if st.session_state.inf_result:
        res = st.session_state.inf_result

        # Auto-inspección
        if st.session_state.inf_inspection:
            st.markdown('<div class="visual-section-title">🔍 Auto-Inspección</div>',
                        unsafe_allow_html=True)
            render_panel_inspeccion(st.session_state.inf_inspection)

        # Galería
        if st.session_state.inf_images:
            render_galeria_imagenes(st.session_state.inf_images, "inf")

        # Vista previa
        st.markdown('<div class="visual-section-title">👁️ Vista previa</div>', unsafe_allow_html=True)
        components.html(res["html"], height=900, scrolling=True)

        # Descargas
        st.markdown('<div class="visual-section-title">📥 Descargar</div>', unsafe_allow_html=True)
        col_d1, col_d2 = st.columns(2)
        with col_d1:
            st.download_button(
                label="📥 Descargar Infografía (.html)",
                data=res["html"].encode("utf-8"),
                file_name=f"Infografia_{ia.sanear_nombre_archivo(res['tema'])}.html",
                mime="text/html", type="primary", width="stretch")
        with col_d2:
            if st.button("🗑️ Descartar infografía", width="stretch", key="discard_inf"):
                for k in ["inf_result", "inf_images", "inf_inspection"]:
                    st.session_state[k] = None
                st.rerun()
        st.caption("💡 Abre el .html en cualquier navegador. Para obtener PDF: "
                   "Imprimir → Guardar como PDF.")